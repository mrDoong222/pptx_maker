import pandas as pd
from pptx import Presentation

# Load the CSV file into a DataFrame
df = pd.read_csv('fiction_class_grid.csv')

# Create a Presentation object
prs = Presentation()  # Call the Presentation class to create an instance

# Loop through each row in the DataFrame
for index, row in df.iterrows():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title_text = str(int(row.iloc[0])) if pd.notna(row.iloc[0]) else 'Student Data'
    title.text = title_text  # Set the title to the student SSID

    content = slide.placeholders[1]
    content_text = str(row.iloc[1]) if pd.notna(row.iloc[1]) else "Name and SSID"
    content.text = content_text  # Set the content for the slide

# Save the presentation
prs.save('new_ppt.pptx')
print(f"Presentation '{'new'}' created successfully.")

